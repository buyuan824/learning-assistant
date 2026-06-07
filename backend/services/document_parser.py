import PyPDF2
from pptx import Presentation
import markdown
import re

class DocumentParser:
    @staticmethod
    def parse_pdf(file_path):
        """解析PDF文件"""
        text = ""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
        except Exception as e:
            raise Exception(f"PDF解析失败: {str(e)}")
        return text
    
    @staticmethod
    def parse_pptx(file_path):
        """解析PPTX文件"""
        text = ""
        try:
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            raise Exception(f"PPT解析失败: {str(e)}")
        return text
    
    @staticmethod
    def parse_markdown(file_path):
        """解析Markdown文件"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                md_content = file.read()
                # 转换为纯文本（简单去除markdown标记）
                text = re.sub(r'[#*_`\[\]]', '', md_content)
                return text
        except Exception as e:
            raise Exception(f"Markdown解析失败: {str(e)}")
    
    @staticmethod
    def parse_file(file_path, file_type):
        """根据文件类型解析文件"""
        parsers = {
            'pdf': DocumentParser.parse_pdf,
            'pptx': DocumentParser.parse_pptx,
            'md': DocumentParser.parse_markdown,
            'txt': lambda fp: open(fp, 'r', encoding='utf-8').read()
        }
        
        parser = parsers.get(file_type.lower())
        if not parser:
            raise Exception(f"不支持的文件类型: {file_type}")
        
        return parser(file_path)
